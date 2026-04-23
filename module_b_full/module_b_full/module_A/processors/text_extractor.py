"""
Extracts plain text from a broad set of educational file formats.
"""

from __future__ import annotations

import json
import re
import subprocess
import zipfile
from pathlib import Path
from typing import Any, Dict


def extract_text(file_path: Path) -> Dict[str, Any]:
    """
    Extract text from a file.
    Returns dict with keys: text, title, metadata.
    """
    ext = file_path.suffix.lower()
    dispatch = {
        ".html": _html,
        ".htm": _html,
        ".xhtml": _html,
        ".xml": _xml,
        ".pdf": _pdf,
        ".docx": _docx,
        ".doc": _doc_legacy,
        ".pptx": _pptx,
        ".ppt": _ppt_legacy,
        ".xlsx": _excel,
        ".xls": _excel,
        ".csv": _csv,
        ".tsv": _tsv,
        ".txt": _plain,
        ".text": _plain,
        ".log": _plain,
        ".md": _plain,
        ".rst": _plain,
        ".rtf": _rtf,
        ".json": _json_text,
        ".jsonl": _jsonl,
        ".yaml": _yaml_text,
        ".yml": _yaml_text,
        ".ini": _plain,
        ".cfg": _plain,
        ".toml": _plain,
        ".odt": _open_document,
        ".ods": _open_document,
        ".odp": _open_document,
        ".epub": _epub,
    }
    fn = dispatch.get(ext, _auto)
    return fn(file_path)


def _make_result(text: str, title: str, fmt: str, **metadata: Any) -> Dict[str, Any]:
    clean = clean_text(text)
    return {
        "text": clean,
        "title": title or (clean.splitlines()[0][:120] if clean else ""),
        "metadata": {"format": fmt, **metadata},
    }


def _error_result(exc: Exception) -> Dict[str, Any]:
    return {"text": "", "title": "", "metadata": {"error": str(exc)}}


def _safe_decode(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _binary_strings(path: Path, min_len: int = 4) -> str:
    data = path.read_bytes()
    chunks = re.findall(rb"[\x20-\x7E\xc0-\xff\t\r\n]{%d,}" % min_len, data)
    seen = set()
    filtered = []
    for chunk in chunks:
        text = clean_text(_safe_decode(chunk))
        if not text or text in seen:
            continue
        if len(text) > 500:
            continue
        letters = sum(ch.isalpha() for ch in text)
        if letters < max(3, len(text) // 10):
            continue
        seen.add(text)
        filtered.append(text)
        if len(filtered) >= 2000:
            break
    return "\n".join(filtered)


def _auto(path: Path) -> Dict[str, Any]:
    try:
        data = path.read_bytes()
        if b"<html" in data.lower() or b"<body" in data.lower():
            return _html(path)
        return _make_result(_safe_decode(data), path.stem, path.suffix.lstrip("."))
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# HTML / XML
# ---------------------------------------------------------------------------

def _html(path: Path) -> Dict[str, Any]:
    try:
        from bs4 import BeautifulSoup

        content = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(content, "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        title = soup.title.get_text(strip=True) if soup.title else path.stem
        main = soup.find(["main", "article"]) or soup.find("body") or soup
        text = main.get_text(separator="\n", strip=True)
        return _make_result(text, title, "html")
    except Exception as exc:
        return _error_result(exc)


def _xml(path: Path) -> Dict[str, Any]:
    try:
        from bs4 import BeautifulSoup

        content = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(content, "xml")
        text = soup.get_text(separator="\n", strip=True)
        return _make_result(text, path.stem, "xml")
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _pdf(path: Path) -> Dict[str, Any]:
    try:
        import pdfplumber

        parts = []
        title = ""
        pages = 0
        with pdfplumber.open(path) as pdf:
            if pdf.metadata:
                title = pdf.metadata.get("Title", "") or ""
            for page in pdf.pages:
                pages += 1
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        row_text = " | ".join(cell or "" for cell in row)
                        if row_text.strip():
                            parts.append(row_text)
                page_text = page.extract_text()
                if page_text:
                    parts.append(page_text)
        return _make_result("\n\n".join(parts), title or path.stem, "pdf", pages=pages)
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# DOCX / DOC
# ---------------------------------------------------------------------------

def _docx(path: Path) -> Dict[str, Any]:
    try:
        from docx import Document

        doc = Document(path)
        title = getattr(doc.core_properties, "title", "") or path.stem
        parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)

        return _make_result("\n".join(parts), title, "docx")
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# PPTX / PPT
# ---------------------------------------------------------------------------

def _pptx(path: Path) -> Dict[str, Any]:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"[Slide {idx}]"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_parts.append(text.strip())
            if len(slide_parts) > 1:
                parts.append("\n".join(slide_parts))
        return _make_result("\n\n".join(parts), path.stem, "pptx", slides=len(prs.slides))
    except Exception as exc:
        return _error_result(exc)


def _run_powershell(script: str) -> str:
    result = subprocess.run(
        ["powershell", "-NoProfile", "-Command", script],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="ignore",
        timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or result.stdout.strip() or "PowerShell failed")
    return result.stdout.strip()


def _doc_legacy(path: Path) -> Dict[str, Any]:
    try:
        file_arg = str(path.resolve()).replace("'", "''")
        script = f"""
$word = $null
$doc = $null
try {{
    $word = New-Object -ComObject Word.Application
    $word.Visible = $false
    $doc = $word.Documents.Open('{file_arg}', $false, $true)
    [Console]::OutputEncoding = [System.Text.Encoding]::UTF8
    $doc.Content.Text
}} finally {{
    if ($doc -ne $null) {{ $doc.Close([ref]0) }}
    if ($word -ne $null) {{ $word.Quit() }}
}}
"""
        text = _run_powershell(script)
        if text:
            return _make_result(text, path.stem, "doc")
    except Exception:
        pass
    return _legacy_binary_text(path)


def _ppt_legacy(path: Path) -> Dict[str, Any]:
    try:
        file_arg = str(path.resolve()).replace("'", "''")
        script = f"""
$ppt = $null
$pres = $null
try {{
    $ppt = New-Object -ComObject PowerPoint.Application
    $pres = $ppt.Presentations.Open('{file_arg}', $false, $false, $false)
    $parts = New-Object System.Collections.Generic.List[string]
    foreach ($slide in $pres.Slides) {{
        $parts.Add("[Slide $($slide.SlideIndex)]")
        foreach ($shape in $slide.Shapes) {{
            try {{
                if ($shape.HasTextFrame -and $shape.TextFrame.HasText) {{
                    $text = $shape.TextFrame.TextRange.Text
                    if ($text) {{ $parts.Add($text) }}
                }}
            }} catch {{}}
        }}
    }}
    [Console]::OutputEncoding = [System.Text.Encoding]::UTF8
    $parts -join [Environment]::NewLine
}} finally {{
    if ($pres -ne $null) {{ $pres.Close() }}
    if ($ppt -ne $null) {{ $ppt.Quit() }}
}}
"""
        text = _run_powershell(script)
        if text:
            return _make_result(text, path.stem, "ppt")
    except Exception:
        pass
    return _legacy_binary_text(path)


def _legacy_binary_text(path: Path) -> Dict[str, Any]:
    try:
        text = _binary_strings(path)
        return _make_result(text, path.stem, path.suffix.lstrip("."))
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def _excel(path: Path) -> Dict[str, Any]:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        sheet_count = len(wb.sheetnames)
        for sheet in wb.worksheets:
            sheet_parts = [f"[Sheet: {sheet.title}]"]
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(cells).strip()
                if row_text.replace("|", "").strip():
                    sheet_parts.append(row_text)
            if len(sheet_parts) > 1:
                parts.append("\n".join(sheet_parts))
        wb.close()
        return _make_result("\n\n".join(parts), path.stem, path.suffix.lstrip("."), sheets=sheet_count)
    except Exception:
        try:
            import pandas as pd

            xls = pd.read_excel(path, sheet_name=None)
            parts = []
            for sheet_name, df in xls.items():
                parts.append(f"[Sheet: {sheet_name}]")
                parts.append(df.to_string(index=False))
            return _make_result("\n\n".join(parts), path.stem, path.suffix.lstrip("."), sheets=len(xls))
        except Exception as exc:
            return _error_result(exc)


# ---------------------------------------------------------------------------
# CSV / TSV
# ---------------------------------------------------------------------------

def _csv(path: Path) -> Dict[str, Any]:
    try:
        import csv

        rows = []
        with open(path, newline="", encoding="utf-8-sig", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        return _make_result("\n".join(rows), path.stem, "csv", rows=len(rows))
    except Exception as exc:
        return _error_result(exc)


def _tsv(path: Path) -> Dict[str, Any]:
    try:
        import csv

        rows = []
        with open(path, newline="", encoding="utf-8-sig", errors="ignore") as f:
            reader = csv.reader(f, delimiter="\t")
            for row in reader:
                rows.append(" | ".join(row))
        return _make_result("\n".join(rows), path.stem, "tsv", rows=len(rows))
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# Structured text
# ---------------------------------------------------------------------------

def _json_text(path: Path) -> Dict[str, Any]:
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        text = json.dumps(data, ensure_ascii=False, indent=2)
        return _make_result(text, path.stem, "json")
    except Exception as exc:
        return _error_result(exc)


def _jsonl(path: Path) -> Dict[str, Any]:
    try:
        rows = []
        for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
            line = line.strip()
            if not line:
                continue
            try:
                rows.append(json.dumps(json.loads(line), ensure_ascii=False))
            except json.JSONDecodeError:
                rows.append(line)
        return _make_result("\n".join(rows), path.stem, "jsonl", rows=len(rows))
    except Exception as exc:
        return _error_result(exc)


def _yaml_text(path: Path) -> Dict[str, Any]:
    try:
        import yaml

        data = yaml.safe_load(path.read_text(encoding="utf-8", errors="ignore"))
        text = yaml.safe_dump(data, allow_unicode=True, sort_keys=False)
        return _make_result(text, path.stem, "yaml")
    except Exception:
        return _plain(path)


def _rtf(path: Path) -> Dict[str, Any]:
    try:
        from striprtf.striprtf import rtf_to_text

        content = path.read_text(encoding="utf-8", errors="ignore")
        return _make_result(rtf_to_text(content), path.stem, "rtf")
    except Exception:
        try:
            content = path.read_text(encoding="utf-8", errors="ignore")
            text = re.sub(r"\\[a-z]+\d* ?", " ", content)
            text = re.sub(r"[{}]", " ", text)
            return _make_result(text, path.stem, "rtf")
        except Exception as exc:
            return _error_result(exc)


# ---------------------------------------------------------------------------
# OpenDocument / EPUB
# ---------------------------------------------------------------------------

def _open_document(path: Path) -> Dict[str, Any]:
    try:
        from bs4 import BeautifulSoup

        with zipfile.ZipFile(path) as zf:
            content = zf.read("content.xml")
        soup = BeautifulSoup(content, "xml")
        text = soup.get_text(separator="\n", strip=True)
        return _make_result(text, path.stem, path.suffix.lstrip("."))
    except Exception as exc:
        return _error_result(exc)


def _epub(path: Path) -> Dict[str, Any]:
    try:
        from bs4 import BeautifulSoup

        parts = []
        with zipfile.ZipFile(path) as zf:
            names = [name for name in zf.namelist() if name.endswith((".html", ".htm", ".xhtml"))]
            for name in names:
                soup = BeautifulSoup(zf.read(name), "lxml")
                parts.append(soup.get_text(separator="\n", strip=True))
        return _make_result("\n\n".join(parts), path.stem, "epub")
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------

def _plain(path: Path) -> Dict[str, Any]:
    try:
        content = path.read_text(encoding="utf-8", errors="ignore")
        title = content.strip().splitlines()[0][:120] if content.strip() else path.stem
        return _make_result(content, title, path.suffix.lstrip("."))
    except Exception as exc:
        return _error_result(exc)


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def clean_text(text: str) -> str:
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def word_count(text: str) -> int:
    return len(text.split()) if text else 0
